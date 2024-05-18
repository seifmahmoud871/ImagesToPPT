from pptx import Presentation
from pptx.util import Inches
import os

prs = Presentation()

# add file path
img_dir = ''

sortedImages = []

for img_name in os.listdir(img_dir):
    if img_name.endswith(('png', 'jpg', 'jpeg', 'bmp', 'gif')):
        img_path = os.path.join(img_dir, img_name)
        sortedImages.append(img_path)

sortedImages.sort()
for img in sortedImages:
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)

        slide.shapes.add_picture(img, Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)


prs.save('output.pptx')
